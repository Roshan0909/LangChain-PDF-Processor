import os
from google import genai
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import io

load_dotenv()
API_KEY = os.getenv("API_KEY")
client = genai.Client(api_key=API_KEY)

def summarize_text(text, summary_type="concise"):
    """Summarize text using Gemini AI"""
    
    if summary_type == "concise":
        prompt = f"""
        Summarize the following text in 3-5 short sentences maximum. 
        Extract only the most critical points. Keep it extremely brief and to the point.
        Make it easy to understand and remember quickly.
        
        Text to summarize:
        {text}
        
        Provide a very short, concise summary (3-5 sentences only):
        """
    elif summary_type == "detailed":
        prompt = f"""
        Provide a detailed summary of the following text.
        Break it down into sections with headings.
        Explain the key concepts thoroughly.
        Include important examples or details.
        
        Text to summarize:
        {text}
        
        Provide a comprehensive summary:
        """
    else:  # bullet points
        prompt = f"""
        Summarize the following text as clear, organized bullet points.
        Group related points under headings.
        Make it easy to scan and memorize.
        
        Text to summarize:
        {text}
        
        Provide bullet-point summary:
        """
    
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt
        )
        return response.text
    except Exception as e:
        return f"Error generating summary: {str(e)}"

def extract_text_from_pdf_file(pdf_file):
    """Extract text from uploaded PDF file"""
    try:
        pdf_reader = PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text
    except Exception as e:
        return None

def extract_text_from_docx_file(docx_file):
    """Extract text from uploaded Word document"""
    try:
        doc = Document(docx_file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + "\t"
                text += "\n"
        return text
    except Exception as e:
        return None

def extract_text_from_pptx_file(pptx_file):
    """Extract text from uploaded PowerPoint presentation"""
    try:
        prs = Presentation(pptx_file)
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return None
