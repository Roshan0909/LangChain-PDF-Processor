import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from google import genai
from dotenv import load_dotenv
import json
import re

load_dotenv()
API_KEY = os.getenv("API_KEY")
client = genai.Client(api_key=API_KEY)

def extract_text_from_pdf(pdf_path, max_pages=20):
    """Extract text from PDF (limit to first max_pages for speed)"""
    try:
        # Try opening with strict=False to handle malformed PDFs
        pdf_reader = PdfReader(pdf_path, strict=False)
        text = ""
        num_pages = min(len(pdf_reader.pages), max_pages)
        
        for i in range(num_pages):
            try:
                page_text = pdf_reader.pages[i].extract_text()
                if page_text:
                    text += page_text + "\n"
            except Exception as page_error:
                print(f"Warning: Could not extract text from page {i+1}: {page_error}")
                continue
        
        if not text.strip():
            print("Warning: No text could be extracted from the PDF")
            return None
            
        return text
    except Exception as e:
        print(f"Error extracting PDF text: {e}")
        if "EOF marker not found" in str(e):
            print("The PDF file appears to be corrupted or incomplete. Please try re-uploading the file.")
        return None

def extract_text_from_docx(docx_path):
    """Extract text from Word document"""
    try:
        doc = Document(docx_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + "\t"
                text += "\n"
        if not text.strip():
            print("Warning: No text could be extracted from the Word document")
            return None
        return text
    except Exception as e:
        print(f"Error extracting Word document text: {e}")
        return None

def extract_text_from_pptx(pptx_path):
    """Extract text from PowerPoint presentation"""
    try:
        prs = Presentation(pptx_path)
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        if not text.strip():
            print("Warning: No text could be extracted from the PowerPoint")
            return None
        return text
    except Exception as e:
        print(f"Error extracting PowerPoint text: {e}")
        return None

def extract_text_from_file(file_path, max_pages=20):
    """Extract text from PDF, Word, or PowerPoint file"""
    try:
        # Check if file exists
        if not os.path.exists(file_path):
            print(f"Error: File not found at {file_path}")
            return None
        
        # Get file extension
        file_ext = file_path.split('.')[-1].lower()
        
        # Extract based on file type
        if file_ext == 'pdf':
            return extract_text_from_pdf(file_path, max_pages)
        elif file_ext == 'docx':
            return extract_text_from_docx(file_path)
        elif file_ext == 'pptx':
            return extract_text_from_pptx(file_path)
        else:
            print(f"Unsupported file type: {file_ext}")
            return None
            
    except Exception as e:
        print(f"Error extracting text from file: {e}")
        print(f"File path: {file_path}")
        return None

def generate_quiz_questions(pdf_text, num_questions=10, topics=None, difficulty='medium'):
    """Generate MCQ questions from PDF text using Gemini AI"""
    
    # Build the prompt with topics if provided
    if topics and topics.strip():
        topic_instruction = f"\n\nIMPORTANT: Generate questions ONLY about the following topics: {topics}\nFocus exclusively on these topics and ignore other content in the text."
    else:
        topic_instruction = ""
    
    # Add difficulty instruction
    difficulty_instructions = {
        'easy': '\n\nDIFFICULTY: EASY - Create simple, straightforward questions testing basic understanding and recall. Avoid complex scenarios or tricky wording.',
        'medium': '\n\nDIFFICULTY: MEDIUM - Create moderately challenging questions that test comprehension and application of concepts.',
        'hard': '\n\nDIFFICULTY: HARD - Create challenging questions requiring deep understanding, analysis, and critical thinking. Include complex scenarios and edge cases.'
    }
    difficulty_instruction = difficulty_instructions.get(difficulty, difficulty_instructions['medium'])
    
    prompt = f"""You are a quiz generator. Based on the following text, generate exactly {num_questions} multiple choice questions.{topic_instruction}{difficulty_instruction}

TEXT:
{pdf_text[:15000]}

Generate {num_questions} multiple choice questions in this EXACT JSON format (no other text):
[
  {{
    "question": "What is...",
    "options": ["Answer 1", "Answer 2", "Answer 3", "Answer 4"],
    "correct_answer": 0
  }}
]

Rules:
- Each question must have exactly 4 options
- correct_answer is the index (0, 1, 2, or 3) of the correct option
- Questions should test key concepts from the text
- Return ONLY valid JSON, nothing else"""
    
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
            config={
                "temperature": 0.7,
                "response_mime_type": "application/json"
            }
        )
        
        response_text = response.text.strip()
        
        # Try to extract JSON if wrapped in markdown
        if "```" in response_text:
            # Find JSON content between code blocks
            lines = response_text.split('\n')
            json_lines = []
            in_code_block = False
            
            for line in lines:
                if line.strip().startswith('```'):
                    in_code_block = not in_code_block
                    continue
                if in_code_block or (line.strip().startswith('[') or line.strip().startswith('{')):
                    json_lines.append(line)
            
            response_text = '\n'.join(json_lines).strip()
        
        # Remove any remaining markdown
        response_text = response_text.replace('```json', '').replace('```', '').strip()
        
        # Parse JSON
        questions = json.loads(response_text)
        
        if not isinstance(questions, list):
            print(f"Response is not a list: {type(questions)}")
            return []
        
        # Validate and clean questions
        validated_questions = []
        for q in questions:
            if isinstance(q, dict) and all(key in q for key in ['question', 'options', 'correct_answer']):
                if isinstance(q['options'], list) and len(q['options']) == 4:
                    if isinstance(q['correct_answer'], int) and 0 <= q['correct_answer'] <= 3:
                        validated_questions.append(q)
        
        if len(validated_questions) == 0:
            print(f"No valid questions found. Raw response: {response_text[:500]}")
        
        return validated_questions[:num_questions]
        
    except json.JSONDecodeError as e:
        print(f"JSON parsing error: {e}")
        print(f"Response text: {response_text[:500]}")
        return []
    except Exception as e:
        print(f"Error generating questions: {e}")
        return []

def generate_quiz_from_pdf(pdf_path, num_questions=10, topics=None, difficulty='medium'):
    """Main function to generate quiz from PDF, Word, or PowerPoint file"""
    
    # Extract text from file (supports PDF, DOCX, PPTX)
    file_text = extract_text_from_file(pdf_path)
    
    if not file_text or len(file_text) < 100:
        file_ext = pdf_path.split('.')[-1].upper()
        return None, f"Could not extract sufficient text from {file_ext} file"
    
    # Generate questions using AI
    questions = generate_quiz_questions(file_text, num_questions, topics, difficulty)
    
    if not questions:
        return None, "Could not generate questions from file content"
    
    return questions, None
