import os
import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
import pandas as pd
import hashlib
import google.generativeai as genai
from docx import Document
from pptx import Presentation

# Load environment variables and configure API
load_dotenv()
API_KEY = os.getenv("API_KEY")
genai.configure(api_key=API_KEY)

# Utility function for hashing files (for caching)
def hash_file(file):
    hasher = hashlib.sha256()
    for chunk in iter(lambda: file.read(4096), b""):
        hasher.update(chunk)
    file.seek(0)  # Reset file pointer
    return hasher.hexdigest()

# Text extraction functions for various file types
def extract_text_from_pdf(file):
    try:
        # Try with strict=False to handle malformed PDFs
        pdf_reader = PdfReader(file, strict=False)
        text = ""
        for page in pdf_reader.pages:
            try:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
            except Exception as page_error:
                print(f"Warning: Could not extract text from a page: {page_error}")
                continue
        
        if not text.strip():
            raise ValueError("No text could be extracted from the PDF")
        return text
    except Exception as e:
        if "EOF marker not found" in str(e):
            raise ValueError("The PDF file is corrupted or incomplete. Please re-upload a valid PDF file.")
        raise ValueError(f"Error reading PDF: {str(e)}")

def extract_text_from_txt(file):
    return file.read().decode("utf-8")

def extract_text_from_csv(file):
    df = pd.read_csv(file)
    return df.to_string()

def extract_text_from_excel(file):
    df = pd.read_excel(file)
    return df.to_string()

def extract_text_from_docx(file):
    """Extract text from Word document (.docx)"""
    try:
        doc = Document(file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + "\t"
                text += "\n"
        if not text.strip():
            raise ValueError("No text could be extracted from the Word document")
        return text
    except Exception as e:
        raise ValueError(f"Error reading Word document: {str(e)}")

def extract_text_from_pptx(file):
    """Extract text from PowerPoint (.pptx)"""
    try:
        prs = Presentation(file)
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        if not text.strip():
            raise ValueError("No text could be extracted from the PowerPoint")
        return text
    except Exception as e:
        raise ValueError(f"Error reading PowerPoint: {str(e)}")

# Function to select the appropriate text extractor based on file type
def extract_text(file, file_type):
    if file_type == "pdf":
        return extract_text_from_pdf(file)
    elif file_type == "txt":
        return extract_text_from_txt(file)
    elif file_type == "csv":
        return extract_text_from_csv(file)
    elif file_type == "xlsx":
        return extract_text_from_excel(file)
    elif file_type in ["docx", "vnd.openxmlformats-officedocument.wordprocessingml.document"]:
        return extract_text_from_docx(file)
    elif file_type in ["pptx", "vnd.openxmlformats-officedocument.presentationml.presentation"]:
        return extract_text_from_pptx(file)
    else:
        raise ValueError("Unsupported file type.")

# Split text into chunks
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=2000, chunk_overlap=200)
    return text_splitter.split_text(text)

# Vector store handling with caching
# Vector store handling with caching
def get_vector_store(text_chunks, file_hash):
    cache_path = f"faiss_index_{file_hash}"
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")  # Ensure embeddings are defined here

    if os.path.exists(cache_path):
        # Load vector store and pass embeddings correctly
        vector_store = FAISS.load_local(cache_path, embeddings, allow_dangerous_deserialization=True)
        return vector_store
    else:
        # Create and save vector store if it doesn't exist
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
        vector_store.save_local(cache_path)
        return vector_store


# QA chain setup
def get_conversational_chain():
    prompt_template = """
    Answer the question in a detailed manner based on the context provided and about the context.Also give answer to questions asked related to this context provided. 
    dont give wrong answer
    Context:\n{context}\nQuestion:\n{question}\n\nAnswer:
    """
    model = ChatGoogleGenerativeAI(model="Gemini 2.0 Flash", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    return load_qa_chain(model, chain_type="stuff", prompt=prompt)

# Main function to handle user queries
def answer_question(user_question, vector_store):
    chain = get_conversational_chain()
    docs = vector_store.similarity_search(user_question)
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response["output_text"]

# Streamlit UI
st.title("File-based Q&A Chatbot")
st.write("Upload a file (PDF, Word, PowerPoint, TXT, CSV, or Excel) and ask questions about its content.")

uploaded_file = st.file_uploader("Upload your file", type=["pdf", "txt", "csv", "xlsx", "docx", "pptx"])

if uploaded_file:
    file_hash = hash_file(uploaded_file)  # Generate a unique hash for caching
    file_type = uploaded_file.type.split('/')[-1]  # Extract file extension type
    
    # Extract and process file content
    try:
        raw_text = extract_text(uploaded_file, file_type)
        text_chunks = get_text_chunks(raw_text)
        vector_store = get_vector_store(text_chunks, file_hash)
        st.success("File processed and ready for Q&A!")

        user_question = st.text_input("Ask a question about the file:")
        if user_question:
            answer = answer_question(user_question, vector_store)
            st.write("Answer:", answer)

    except Exception as e:
        st.error(f"Error processing file: {e}")
